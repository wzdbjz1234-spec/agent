"""真实解析器、提取物和索引的组合测试。"""

from __future__ import annotations

import json
import sqlite3
import stat
from pathlib import Path

import duckdb
import pyarrow as pa
import pyarrow.parquet as pq
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter

from dataharness.domain import FileVersionStatus
from dataharness.idgen import DeterministicIdFactory
from dataharness.projects import ProjectCorpus
from dataharness.providers.workspace import LocalWorkspace
from dataharness.storage import RuntimeConnectionFactory, SqliteRuntimeStore


def _fixtures(root: Path) -> tuple[Path, ...]:
    csv_path = root / "rows.csv"
    csv_path.write_text("name,value\nalpha,1\n", encoding="utf-8")
    json_path = root / "rows.json"
    json_path.write_text(json.dumps([{"name": "beta", "value": 2}]), encoding="utf-8")
    markdown = root / "notes.md"
    markdown.write_text("# gamma heading", encoding="utf-8")
    text = root / "plain.txt"
    text.write_text("delta paragraph", encoding="utf-8")

    parquet_path = root / "rows.parquet"
    pq.write_table(pa.table({"name": ["epsilon"], "value": [3]}), parquet_path)

    xlsx_path = root / "book.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    assert sheet is not None
    sheet.append(["name", "value"])
    sheet.append(["zeta", 4])
    workbook.save(xlsx_path)

    pdf_path = root / "blank.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=72, height=72)
    with pdf_path.open("wb") as handle:
        writer.write(handle)

    docx_path = root / "document.docx"
    document = Document()
    document.add_paragraph("eta paragraph")
    document.save(str(docx_path))

    pptx_path = root / "slides.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.add_textbox(Inches(0), Inches(0), Inches(2), Inches(1)).text = "theta slide"
    presentation.save(str(pptx_path))

    sqlite_path = root / "snapshot.sqlite"
    connection = sqlite3.connect(sqlite_path)
    connection.execute("CREATE TABLE facts(name TEXT)")
    connection.execute("INSERT INTO facts VALUES ('iota')")
    connection.commit()
    connection.close()

    duckdb_path = root / "snapshot.duckdb"
    duck = duckdb.connect(str(duckdb_path))
    duck.execute("CREATE TABLE facts AS SELECT 'kappa' AS name")
    duck.close()
    return (
        csv_path,
        json_path,
        markdown,
        text,
        parquet_path,
        xlsx_path,
        pdf_path,
        docx_path,
        pptx_path,
        sqlite_path,
        duckdb_path,
    )


def test_all_v1_formats_import_to_versioned_extraction_and_index(tmp_path: Path) -> None:
    store = SqliteRuntimeStore(RuntimeConnectionFactory(tmp_path / "runtime.db"))
    workspace = LocalWorkspace(tmp_path / "projects")
    corpus = ProjectCorpus(store, workspace, id_factory=DeterministicIdFactory())
    project = corpus.create_project("formats")

    versions = [corpus.import_file(project.id, path) for path in _fixtures(tmp_path)]
    assert {version.status for version in versions} == {FileVersionStatus.READY}
    for version in versions:
        extracted = workspace.extracted_path(project.id, version.id)
        payload = json.loads(extracted.read_text("utf-8"))
        assert payload["source_hash"] == version.content_hash
        assert payload["extractor_version"] == "builtin-v1"

    snapshot = corpus.create_snapshot(project.id)
    assert corpus.search(snapshot.id, "alpha")[0].file_version_id == versions[0].id
    assert corpus.search(snapshot.id, "theta")[0].locator == {"slide": 1}

    # 提取物与索引是可重建派生物；删除后从不可变 source 恢复，不改变旧 Snapshot。
    extracted = workspace.extracted_path(project.id, versions[0].id)
    extracted.chmod(stat.S_IWRITE)
    extracted.unlink()
    workspace.index_path(project.id).unlink()
    corpus.rebuild_derived(versions[0].id)
    assert corpus.search(snapshot.id, "alpha")[0].file_version_id == versions[0].id


def test_corrupt_supported_file_is_failed_not_ready(tmp_path: Path) -> None:
    store = SqliteRuntimeStore(RuntimeConnectionFactory(tmp_path / "runtime.db"))
    corpus = ProjectCorpus(
        store,
        LocalWorkspace(tmp_path / "projects"),
        id_factory=DeterministicIdFactory(),
    )
    project = corpus.create_project("corrupt")
    corrupt = tmp_path / "broken.json"
    corrupt.write_text("{not valid JSON", encoding="utf-8")
    assert corpus.import_file(project.id, corrupt).status == FileVersionStatus.FAILED
