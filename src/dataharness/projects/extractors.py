"""内置文档提取器 Adapter 与真实格式嗅探。"""

from __future__ import annotations

import csv
import json
import sqlite3
import zipfile
from pathlib import Path
from typing import Protocol

import duckdb
import pyarrow.parquet as parquet
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

from dataharness.domain import ContentHash

from .models import ExtractedDocument, TextChunk

EXTRACTOR_VERSION = "builtin-v1"
SUPPORTED_MEDIA_TYPES = frozenset(
    {
        "text/csv",
        "application/vnd.apache.parquet",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/json",
        "application/pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "text/markdown",
        "text/plain",
        "application/vnd.sqlite3",
        "application/vnd.duckdb",
    }
)


class UnsupportedFormatError(ValueError):
    """真实内容格式不在 V1 支持列表中。"""


class DocumentExtractor(Protocol):
    """内部提取器协议；第三方对象不会出现在返回值中。"""

    def __call__(
        self, path: Path, source_hash: ContentHash, media_type: str
    ) -> ExtractedDocument: ...


def sniff_media_type(path: Path) -> str:
    """优先按内容魔数识别，文本家族才参考扩展名消除歧义。"""
    header = path.read_bytes()[:16]
    suffix = path.suffix.casefold()
    if header.startswith(b"%PDF-"):
        return "application/pdf"
    if header.startswith(b"PAR1"):
        return "application/vnd.apache.parquet"
    if header.startswith(b"SQLite format 3\x00"):
        return "application/vnd.sqlite3"
    if len(header) >= 12 and header[8:12] == b"DUCK":
        return "application/vnd.duckdb"
    if zipfile.is_zipfile(path):
        with zipfile.ZipFile(path) as archive:
            names = set(archive.namelist())
        if "xl/workbook.xml" in names:
            return "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        if "word/document.xml" in names:
            return "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        if "ppt/presentation.xml" in names:
            return "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        raise UnsupportedFormatError("未知 ZIP 容器不是受支持的 Office 文档")
    if suffix == ".json":
        json.loads(path.read_text(encoding="utf-8-sig"))
        return "application/json"
    if suffix == ".csv":
        path.read_text(encoding="utf-8-sig")
        return "text/csv"
    if suffix in {".md", ".markdown"}:
        path.read_text(encoding="utf-8-sig")
        return "text/markdown"
    if suffix in {".txt", ".text"}:
        path.read_text(encoding="utf-8-sig")
        return "text/plain"
    raise UnsupportedFormatError(f"无法识别真实格式：{suffix or '<无扩展名>'}")


def _document(
    source_hash: ContentHash, media_type: str, chunks: list[TextChunk]
) -> ExtractedDocument:
    return ExtractedDocument(
        source_hash=source_hash,
        media_type=media_type,
        extractor_version=EXTRACTOR_VERSION,
        chunks=tuple(chunks),
    )


def _extract_text(path: Path, source_hash: ContentHash, media_type: str) -> ExtractedDocument:
    chunks = [
        TextChunk(text=line, locator={"paragraph": number})
        for number, line in enumerate(path.read_text("utf-8-sig").splitlines(), 1)
        if line.strip()
    ]
    return _document(source_hash, media_type, chunks)


def _extract_csv(path: Path, source_hash: ContentHash, media_type: str) -> ExtractedDocument:
    chunks: list[TextChunk] = []
    with path.open("r", encoding="utf-8-sig", newline="") as handle:
        rows = list(csv.reader(handle))
    schema = rows[0] if rows else []
    for number, row in enumerate(rows[1:] if rows else [], 2):
        chunks.append(
            TextChunk(
                text=" | ".join(
                    f"{schema[i] if i < len(schema) else i}: {value}" for i, value in enumerate(row)
                ),
                locator={"row_start": number, "row_end": number},
                metadata={"columns": schema},
            )
        )
    return _document(source_hash, media_type, chunks)


def _extract_json(path: Path, source_hash: ContentHash, media_type: str) -> ExtractedDocument:
    value = json.loads(path.read_text("utf-8-sig"))
    values = value if isinstance(value, list) else [value]
    chunks = [
        TextChunk(
            text=json.dumps(item, ensure_ascii=False, sort_keys=True),
            locator={"item": number},
            metadata={"type": type(item).__name__},
        )
        for number, item in enumerate(values)
    ]
    return _document(source_hash, media_type, chunks)


def _extract_xlsx(path: Path, source_hash: ContentHash, media_type: str) -> ExtractedDocument:
    workbook = load_workbook(path, read_only=True, data_only=True)
    chunks: list[TextChunk] = []
    try:
        for sheet in workbook.worksheets:
            rows = sheet.iter_rows(values_only=True)
            header = tuple(next(rows, ()))
            for number, row in enumerate(rows, 2):
                text = " | ".join(
                    f"{header[i] if i < len(header) and header[i] is not None else i}: {value}"
                    for i, value in enumerate(row)
                    if value is not None
                )
                if text:
                    chunks.append(
                        TextChunk(
                            text=text,
                            locator={
                                "worksheet": sheet.title,
                                "row_start": number,
                                "row_end": number,
                            },
                            metadata={"columns": [str(item) for item in header]},
                        )
                    )
    finally:
        workbook.close()
    return _document(source_hash, media_type, chunks)


def _extract_parquet(path: Path, source_hash: ContentHash, media_type: str) -> ExtractedDocument:
    table = parquet.read_table(path)
    chunks = [
        TextChunk(
            text=json.dumps(row, ensure_ascii=False, default=str, sort_keys=True),
            locator={"row_start": number, "row_end": number},
            metadata={"columns": table.column_names, "schema": str(table.schema)},
        )
        for number, row in enumerate(table.to_pylist(), 1)
    ]
    return _document(source_hash, media_type, chunks)


def _extract_pdf(path: Path, source_hash: ContentHash, media_type: str) -> ExtractedDocument:
    reader = PdfReader(path)
    chunks = [
        TextChunk(text=text, locator={"page": number})
        for number, page in enumerate(reader.pages, 1)
        if (text := (page.extract_text() or "").strip())
    ]
    return _document(source_hash, media_type, chunks)


def _extract_docx(path: Path, source_hash: ContentHash, media_type: str) -> ExtractedDocument:
    document = Document(str(path))
    chunks = [
        TextChunk(text=p.text, locator={"paragraph": number})
        for number, p in enumerate(document.paragraphs, 1)
        if p.text.strip()
    ]
    for table_number, table in enumerate(document.tables, 1):
        for row_number, row in enumerate(table.rows, 1):
            chunks.append(
                TextChunk(
                    text=" | ".join(cell.text for cell in row.cells),
                    locator={"table": table_number, "row_start": row_number, "row_end": row_number},
                )
            )
    return _document(source_hash, media_type, chunks)


def _extract_pptx(path: Path, source_hash: ContentHash, media_type: str) -> ExtractedDocument:
    presentation = Presentation(str(path))
    chunks: list[TextChunk] = []
    for slide_number, slide in enumerate(presentation.slides, 1):
        texts = (str(getattr(shape, "text", "")) for shape in slide.shapes)
        text = "\n".join(value for value in texts if value.strip())
        if text:
            chunks.append(TextChunk(text=text, locator={"slide": slide_number}))
    return _document(source_hash, media_type, chunks)


def _extract_sqlite(path: Path, source_hash: ContentHash, media_type: str) -> ExtractedDocument:
    connection = sqlite3.connect(f"file:{path.as_posix()}?mode=ro", uri=True)
    try:
        tables = connection.execute(
            "SELECT name, sql FROM sqlite_master WHERE type='table' "
            "AND name NOT LIKE 'sqlite_%' ORDER BY name"
        ).fetchall()
        chunks = []
        for name, schema in tables:
            count = connection.execute(
                f'SELECT COUNT(*) FROM "{name.replace(chr(34), chr(34) * 2)}"'
            ).fetchone()[0]
            chunks.append(
                TextChunk(
                    text=f"table {name}; rows={count}; schema={schema}",
                    locator={"table": name},
                    metadata={"row_count": count},
                )
            )
        return _document(source_hash, media_type, chunks)
    finally:
        connection.close()


def _extract_duckdb(path: Path, source_hash: ContentHash, media_type: str) -> ExtractedDocument:
    connection = duckdb.connect(str(path), read_only=True)
    try:
        tables = connection.execute("SHOW ALL TABLES").fetchall()
        chunks = [
            TextChunk(text=" | ".join(map(str, row)), locator={"table": str(row[2])})
            for row in tables
        ]
        return _document(source_hash, media_type, chunks)
    finally:
        connection.close()


_EXTRACTORS: dict[str, DocumentExtractor] = {
    "text/plain": _extract_text,
    "text/markdown": _extract_text,
    "text/csv": _extract_csv,
    "application/json": _extract_json,
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": _extract_xlsx,
    "application/vnd.apache.parquet": _extract_parquet,
    "application/pdf": _extract_pdf,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": _extract_docx,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": _extract_pptx,
    "application/vnd.sqlite3": _extract_sqlite,
    "application/vnd.duckdb": _extract_duckdb,
}


def extract_document(path: Path, source_hash: ContentHash, media_type: str) -> ExtractedDocument:
    """按稳定媒体类型分派内部提取器。"""
    extractor = _EXTRACTORS.get(media_type)
    if extractor is None:
        raise UnsupportedFormatError(media_type)
    return extractor(path, source_hash, media_type)
